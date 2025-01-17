import os
import io
import time
import threading
import uuid
import json
import re
import html
import requests
from flask import Flask, request, jsonify, send_file, url_for
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

POTX_URL = "https://files.assistantos.de/powerpoint/Master.potx"
generated_files = {}  # Speichert Dateien und deren Erstellungszeit
TTL = 7 * 24 * 60 * 60  # Lebensdauer der Datei in Sekunden (7 Tage)

def cleanup_files():
    """Entfernt abgelaufene Dateien."""
    while True:
        time.sleep(60)  # Überprüfung alle 60 Sekunden
        now = time.time()
        expired_files = [file_id for file_id, (file_path, timestamp) in generated_files.items() if now - timestamp > TTL]
        for file_id in expired_files:
            try:
                os.remove(file_path)  # Löscht die Datei
                del generated_files[file_id]  # Entfernt Eintrag aus dem Dictionary
                print(f"Datei {file_id} wurde gelöscht.")
            except Exception as e:
                print(f"Fehler beim Löschen der Datei {file_id}: {e}")

# Startet den Aufräum-Thread
threading.Thread(target=cleanup_files, daemon=True).start()

def escape_text(text: str) -> str:
    """Entfernt Markdown und HTML-Artefakte."""
    if not text:
        return ""
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'\\\\n', r'\n', text)
    text = text.replace(r'\n', '\n')
    text = re.sub(r'^\s*\*\s+', '- ', text, flags=re.MULTILINE)
    text = html.unescape(text)
    text = text.replace(r'\"', '"')
    return text

@app.route('/generate_pptx', methods=['POST', 'OPTIONS'])
def generate_pptx():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'OK'}), 200

    try:
        request_body = request.get_data(as_text=True)
        data = json.loads(request_body)
        slides_data = data.get('slides', {})

        potx_response = requests.get(POTX_URL)
        potx_response.raise_for_status()
        prs = Presentation(io.BytesIO(potx_response.content))

        # Title Slide
        if 'title_slide' in slides_data:
            title_slide_data = slides_data['title_slide']
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)
            title_slide.shapes.title.text = escape_text(title_slide_data.get('title', 'Presentation Title'))
            subtitle_placeholder = title_slide.placeholders[1]
            subtitle_placeholder.text = escape_text(title_slide_data.get('subtitle', ''))

        # Table of Contents
        if 'table_of_contents' in slides_data:
            toc_data = slides_data['table_of_contents']
            toc_slide_layout = prs.slide_layouts[1]
            toc_slide = prs.slides.add_slide(toc_slide_layout)
            toc_slide.shapes.title.text = escape_text(toc_data.get('title', 'Inhaltsverzeichnis'))
            for shape in toc_slide.shapes:
                if shape.has_text_frame and not shape.text_frame.text.strip():
                    text_frame = shape.text_frame
                    text_frame.text = "\n".join(f"{escape_text(entry)}" for entry in toc_data.get('entries', []))
                    break

        # Content Slides
        for slide_data in slides_data.get('content_slides', []):
            slide_layout = prs.slide_layouts[2] if any(k in slide_data for k in ['chart_data', 'images', 'table_data']) else prs.slide_layouts[3]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = escape_text(slide_data.get('title', 'Untitled Slide'))

            # Body Text
            if 'body' in slide_data:
                body_placeholder = None
                for shape in slide.placeholders:
                    if shape.has_text_frame and not shape.text_frame.text.strip():
                        body_placeholder = shape
                        break
                if body_placeholder:
                    text_frame = body_placeholder.text_frame
                    text_frame.text = escape_text(slide_data['body'])
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP

            # Bilder
            if 'images' in slide_data:
                for img_data in slide_data['images']:
                    try:
                        headers = {"User-Agent": "Powerpoint_Generator_bot/1.0"}
                        response = requests.get(img_data['url'], headers=headers, stream=True)
                        response.raise_for_status()
                        image_stream = io.BytesIO(response.content)

                        left = Inches(img_data.get('left', 1))
                        top = Inches(img_data.get('top', 1))
                        width = Inches(img_data.get('width', 3))
                        height = Inches(img_data.get('height', 2))
                        slide.shapes.add_picture(image_stream, left, top, width, height)
                    except Exception as e:
                        print(f"Error adding image: {e}")

            # Tabellen
            if 'table_data' in slide_data:
                table_data = slide_data['table_data']
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0

                table_position = slide_data.get('table_position', {})
                left = Inches(table_position.get('left', 1))
                top = Inches(table_position.get('top', 3))
                width = Inches(table_position.get('width', 8))
                height = Inches(table_position.get('height', 2))

                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                table = table_shape.table

                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_data in enumerate(row_data):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_data)

            # Diagramme
            if 'chart_data' in slide_data:
                chart_data_input = slide_data['chart_data']
                try:
                    chart_data_obj = CategoryChartData()
                    chart_data_obj.categories = chart_data_input.get('categories', [])
                    for series in chart_data_input.get('series', []):
                        chart_data_obj.add_series(series.get('name', ''), series.get('values', []))

                    chart_position = chart_data_input.get('chart_position', {})
                    x = Inches(chart_position.get('left', 1))
                    y = Inches(chart_position.get('top', 3))
                    cx = Inches(chart_position.get('width', 6))
                    cy = Inches(chart_position.get('height', 4))

                    chart_type_str = chart_data_input.get('type', 'COLUMN_CLUSTERED')
                    chart_type = getattr(XL_CHART_TYPE, chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

                    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data_obj).chart

                    chart.has_legend = chart_data_input.get('has_legend', True)
                    if chart_data_input.get('title'):
                        chart.chart_title.has_text_frame = True
                        chart.chart_title.text_frame.text = chart_data_input['title']

                except Exception as e:
                    print(f"Error adding chart: {e}")

        # Closing Slide
        if 'closing_slide' in slides_data:
            closing_slide_data = slides_data['closing_slide']
            closing_slide_layout = prs.slide_layouts[4]
            closing_slide = prs.slides.add_slide(closing_slide_layout)
            closing_slide.shapes.title.text = escape_text(closing_slide_data.get('title', 'Vielen Dank'))
            for shape in closing_slide.shapes:
                if shape.has_text_frame and not shape.text_frame.text.strip():
                    text_frame = shape.text_frame
                    text_frame.text = escape_text(closing_slide_data.get('body', ''))
                    break

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        file_id = str(uuid.uuid4())
        file_path = f"/tmp/{file_id}.pptx"
        with open(file_path, 'wb') as f:
            f.write(pptx_buffer.getvalue())

        # Speichere Dateipfad und Erstellungszeit
        generated_files[file_id] = (file_path, time.time())
        download_link = url_for('download_file', file_id=file_id, _external=True)
        return jsonify({'download_link': download_link})

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download/<file_id>')
def download_file(file_id):
    if file_id in generated_files:
        file_path, _ = generated_files[file_id]
        return send_file(
            file_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            download_name='generated_presentation.pptx',
            as_attachment=True
        )
    return "File not found", 404
